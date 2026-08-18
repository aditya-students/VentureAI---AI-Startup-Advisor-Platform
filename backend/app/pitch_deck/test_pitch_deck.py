"""
Unit tests for AI Pitch Deck Generator.

Tests:
1. Context prerequisite checker & missing upstream feature handling
2. Low Validation Mode flag trigger (< 50 score)
3. Red Pen Auditor contradiction detection & auto-rewriting
4. Fallback slide generators & Schema validation for all 13 slides
"""

import unittest
from unittest.mock import MagicMock

from app.pitch_deck.context import get_prerequisites_status, build_pitch_deck_context
from app.pitch_deck.graph.nodes import (
    run_pitch_deck_audit,
    _fallback_group_a,
    _fallback_group_b,
    _fallback_group_c,
)
from app.pitch_deck.schemas import (
    SlideSchema,
    PitchDeckAuditReport,
)


class TestPitchDeckContext(unittest.TestCase):
    def test_prerequisites_status(self):
        db = MagicMock()
        startup_mock = MagicMock()
        startup_mock.id = 1

        def mock_query(model):
            q = MagicMock()
            if model.__name__ == "Startup":
                q.filter().first.return_value = startup_mock
            else:
                q.filter().order_by().first.return_value = None
            return q

        db.query.side_effect = mock_query

        res = get_prerequisites_status(db, 1)
        self.assertTrue(res["workspace_exists"])
        self.assertFalse(res["validation_exists"])
        self.assertFalse(res["bmc_exists"])
        self.assertFalse(res["business_plan_exists"])
        self.assertIn("Your Idea Validation, Business Model Canvas, Business Plan has not been generated yet", res["missing_message"])


class TestPitchDeckRedPenAuditor(unittest.TestCase):
    def test_moat_contradiction_auto_rewrite(self):
        context = {
            "startup_data": {
                "name": "AcmeAI",
                "problem": "Manual data entry",
                "solution": "Automated data entry",
                "target_market": "Local Businesses",
            },
            "validation_data": {
                "final_validation_score": 45,
                "dimension_scores": {"moat": 20},
            },
            "validation_score": 45.0,
            "is_validation_mode": True,
        }

        slides = [
            {"slide_number": 10, "slide_type": "defensive_moat", "title": "Defensive Moat", "subtitle": "Moat", "content": "We have an unbeatable proprietary AI patent technology moat.", "key_points": ["Patented tech"], "visual_type": "grid_cards", "visual_data": {}, "icon_names": ["lock"]}
        ]

        audit = run_pitch_deck_audit(context, slides)
        self.assertTrue(any(w["category"] == "Moat Contradiction" for w in audit["warnings"]))
        # Verify auto-rewrite
        self.assertNotIn("unbeatable proprietary", slides[0]["content"].lower())
        self.assertIn("execution speed", slides[0]["content"].lower())

    def test_fake_traction_detection(self):
        context = {
            "startup_data": {"name": "TractionApp", "target_market": "Consumers"},
            "validation_data": {"lofa": "Users will sign up"},
            "validation_score": 60.0,
            "is_validation_mode": False,
        }

        slides = [
            {"slide_number": 11, "slide_type": "validation_traction", "title": "Traction", "subtitle": "Metrics", "content": "We have 10,000 active users and $100k MRR growing 50% MoM.", "key_points": ["$100k MRR"], "visual_type": "blueprint_list", "visual_data": {}, "icon_names": ["chart"]}
        ]

        audit = run_pitch_deck_audit(context, slides)
        self.assertTrue(any(w["category"] == "Fake Traction" for w in audit["warnings"]))
        self.assertIn("Pre-revenue", slides[0]["content"])


class TestPitchDeckFallbacks(unittest.TestCase):
    def test_all_13_slides_schema_validation(self):
        startup = {
            "name": "VentureApp",
            "tagline": "AI Startup Advisor",
            "problem": "Founders struggle with business planning",
            "solution": "Automated AI business advisor platform",
            "industry": "EdTech / SaaS",
            "target_market": "Early stage founders",
            "stage": "Idea",
            "founder": {"name": "Alice Founder"},
        }
        val = {
            "lofa": "Founders will pay for automated pitch deck generation",
            "dimension_scores": {"problem": 60, "buyer": 65, "market": 70, "moat": 25, "feasibility": 80},
            "mom_test_questions": ["How do you create your pitch deck today?"],
            "kill_threshold": "Fewer than 10 founder interviews express intent to pay"
        }
        bmc = {"canvas_blocks": {"channels": ["Organic SEO", "Founder Communities"], "revenue_streams": ["SaaS Subscription"]}}
        bp = {}

        slides_a = _fallback_group_a(startup, val, bp, 55.0, False)
        slides_b = _fallback_group_b(startup, val, bmc, bp, 25.0)  # Moat score 25
        slides_c = _fallback_group_c(startup, val, bmc, bp, 55.0, False)

        all_slides = slides_a + slides_b + slides_c
        self.assertEqual(len(all_slides), 13)

        # Validate each slide using SlideSchema
        for slide_dict in all_slides:
            slide_obj = SlideSchema(**slide_dict)
            self.assertGreaterEqual(slide_obj.slide_number, 1)
            self.assertLessEqual(slide_obj.slide_number, 13)

        # Verify Slide 10 (Moat < 30) warning
        slide10 = next(s for s in all_slides if s["slide_number"] == 10)
        self.assertTrue(any("Moat score <30" in w for w in slide10["warnings"]))

        # Verify Slide 11 (Traction) anti-hallucination wording
        slide11 = next(s for s in all_slides if s["slide_number"] == 11)
        self.assertIn("Pre-revenue", slide11["subtitle"])


class TestPitchDeckSanitizer(unittest.TestCase):
    def test_sanitize_text_unwraps_dict_strings(self):
        from app.pitch_deck.service import _sanitize_text_py, _parse_hero_stat_py
        raw = "Acquisition channels: {'items': ['Direct digital outreach & content inbound', 'Targeted professional social channels'], 'risk_notes': None, 'last_updated': '2026-08-12T18:35:01', 'generated_by_ai': True}."
        sanitized = _sanitize_text_py(raw)
        self.assertNotIn("risk_notes", sanitized)
        self.assertNotIn("generated_by_ai", sanitized)
        self.assertIn("Direct digital outreach & content inbound", sanitized)
        self.assertIn("Targeted professional social channels", sanitized)

    def test_parse_hero_stat_ignores_date_years(self):
        from app.pitch_deck.service import _parse_hero_stat_py
        stat1 = _parse_hero_stat_py("2026-08-12T18:35:01")
        self.assertIsNone(stat1)
        stat2 = _parse_hero_stat_py("$14.2B Annual Fines")
        self.assertIsNotNone(stat2)
        self.assertEqual(stat2["num"], "$14.2B")
        self.assertEqual(stat2["label"], "Annual Fines")


class TestPPTXBuilder(unittest.TestCase):
    def test_pptx_clean_bullet_text_and_extract_metric(self):
        from app.pitch_deck.pptx_builder import clean_bullet_text, extract_metric, build_pitch_deck_pptx

        raw_dict_str = "Acquisition channels: {'items': ['Direct digital outreach & content inbound', 'Targeted professional social channels'], 'risk_notes': None}."
        cleaned = clean_bullet_text(raw_dict_str)
        self.assertNotIn("risk_notes", cleaned)
        self.assertIn("Direct digital outreach & content inbound", cleaned)

        num, label = extract_metric("$14.2B Annual SEC Fines")
        self.assertEqual(num, "$14.2B")
        self.assertEqual(label, "Annual SEC Fines")

        num_date, label_date = extract_metric("2026-08-12T18:35:01")
        self.assertEqual(num_date, "")

    def test_build_pitch_deck_pptx_returns_bytes(self):
        from app.pitch_deck.pptx_builder import build_pitch_deck_pptx
        sample_slides = [
            {"slide_number": 1, "slide_type": "cover", "title": "SentinelLedger AI", "subtitle": "AI Middleware", "content": "Zero trust banking", "key_points": ["$14.2B Fines", "97% False Positives"]},
            {"slide_number": 6, "slide_type": "product_workflow", "title": "Product Workflow", "subtitle": "3-Step Pipeline", "content": "How it works", "key_points": ["Step 1", "Step 2", "Step 3"], "visual_data": {"steps": [{"step": 1, "title": "Ingestion", "desc": "Data ingest"}, {"step": 2, "title": "Analysis", "desc": "Multi-agent"}, {"step": 3, "title": "Output", "desc": "Audit ready"}]}}
        ]
        pptx_bytes = build_pitch_deck_pptx(sample_slides, "SentinelLedger AI")
        self.assertIsInstance(pptx_bytes, bytes)
        self.assertGreater(len(pptx_bytes), 1000)

    def test_pptx_coordinate_budget_and_shape_properties(self):
        import io
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        from app.pitch_deck.pptx_builder import build_pitch_deck_pptx

        sample_slides = [
            {
                "slide_number": 2,
                "slide_type": "problem_statement",
                "title": "Problem Statement",
                "subtitle": "High manual oversight in compliance",
                "content": "Compliance teams waste 15-20 hours per week.",
                "key_points": [
                    "Manual bottleneck in review",
                    "High cost per audit",
                    "Slow turnaround times",
                    "Human error in validation"
                ]
            },
            {
                "slide_number": 6,
                "slide_type": "product_workflow",
                "title": "Product Workflow",
                "subtitle": "3-Step Pipeline",
                "content": "How it works",
                "key_points": ["Step 1", "Step 2", "Step 3"],
                "visual_data": {
                    "steps": [
                        {"step": 1, "title": "Ingestion", "desc": "Data ingest"},
                        {"step": 2, "title": "Analysis", "desc": "Multi-agent"},
                        {"step": 3, "title": "Output", "desc": "Audit ready"}
                    ]
                }
            }
        ]

        pptx_bytes = build_pitch_deck_pptx(sample_slides, "VendorTrust AI")
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Check slide dimensions (16:9 widescreen)
        self.assertAlmostEqual(prs.slide_width.inches, 13.333, places=2)
        self.assertAlmostEqual(prs.slide_height.inches, 7.5, places=2)

        max_allowed_y = Inches(6.6)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                    bottom_y_inches = (shape.top + shape.height) / Inches(1)
                    self.assertLessEqual(
                        bottom_y_inches,
                        6.6,
                        f"Card shape on slide {slide_idx + 1} bottom ({bottom_y_inches:.2f}\") exceeds max allowed Y of 6.6\""
                    )
                    self.assertEqual(shape.auto_shape_type, MSO_SHAPE.ROUNDED_RECTANGLE)



if __name__ == "__main__":
    unittest.main()
