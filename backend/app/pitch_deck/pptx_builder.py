"""
Enterprise PPTX Builder for VentureAI Pitch Decks using python-pptx.
Implements data sanitization, metric extraction, smart text truncation,
16:9 widescreen layout, and dark enterprise card containers.
"""

import io
import json
import re
from typing import List, Dict, Any, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def clean_bullet_text(item: Any, max_len: int = 160) -> str:
    """
    Extracts text if the item is a dict, list, or dirty string.
    Removes trailing/broken quotes and cleans up JSON string artifacts.
    Smartly truncates text at max_len without breaking words.
    """
    if item is None:
        return ""

    if isinstance(item, dict):
        if "items" in item and isinstance(item["items"], list):
            return clean_bullet_text(item["items"], max_len)
        if "title" in item or "desc" in item or "label" in item:
            parts = [item.get("title"), item.get("desc") or item.get("label")]
            return clean_bullet_text(": ".join([str(p) for p in parts if p]), max_len)
        clean_vals = []
        for k, v in item.items():
            if k in ['risk_notes', 'last_updated', 'generated_by_ai', 'modified_by_founder', 'id', 'created_at']:
                continue
            if v:
                clean_vals.append(clean_bullet_text(v, max_len))
        return clean_bullet_text(" • ".join(clean_vals), max_len)

    if isinstance(item, list):
        return clean_bullet_text(" • ".join([clean_bullet_text(it, 100) for it in item if clean_bullet_text(it, 100)]), max_len)

    text = str(item).strip()

    # Unwrap embedded dict strings e.g. Acquisition channels: {'items': [...]}
    if "{" in text and "}" in text:
        def replace_items_dict(match):
            items_str = match.group(1)
            items = [re.sub(r"^['\"\s]+|['\"\s]+$", "", it) for it in items_str.split(",") if it.strip()]
            return ", ".join(items)

        text = re.sub(r"\{[^{}]*['\"]items['\"]\s*:\s*\[([^\]]+)\][^{}]*\}", replace_items_dict, text, flags=re.IGNORECASE)

        if text.startswith("{") and text.endswith("}"):
            try:
                json_str = text.replace("'", '"').replace("None", "null").replace("True", "true").replace("False", "false")
                parsed = json.loads(json_str)
                return clean_bullet_text(parsed, max_len)
            except Exception:
                text = re.sub(r"'risk_notes':\s*None,?", "", text, flags=re.IGNORECASE)
                text = re.sub(r"'last_updated':\s*'[^']+',?", "", text, flags=re.IGNORECASE)
                text = re.sub(r"'generated_by_ai':\s*(True|False),?", "", text, flags=re.IGNORECASE)
                text = re.sub(r"'modified_by_founder':\s*(True|False),?", "", text, flags=re.IGNORECASE)
                text = re.sub(r"['\"]items['\"]:\s*", "", text, flags=re.IGNORECASE)
                text = re.sub(r"[{}'\"]", "", text)

    # Clean residual dict artifacts & quotes
    text = re.sub(r"'risk_notes':\s*None,?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"'last_updated':\s*'[^']+',?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"'generated_by_ai':\s*(True|False),?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"'modified_by_founder':\s*(True|False),?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^['\"]+|['\"]+$", "", text)
    text = re.sub(r"\s+", " ", text).strip()

    # Smart truncation at max_len
    if len(text) > max_len:
        truncated = text[:max_len].rsplit(" ", 1)[0]
        text = truncated.rstrip(".,;:") + "..."

    return text


def extract_metric(text: str) -> Tuple[str, str]:
    """
    Detects if a bullet point contains a metric (e.g. '$14.2B', '97%', '88%', '<30s', '3.0x', '$100k').
    Returns (metric_string, remaining_label_string).
    """
    cleaned = clean_bullet_text(text, max_len=200)
    if not cleaned:
        return ("", "")

    # Reject ISO timestamps or date fields
    if re.match(r"^\d{4}-\d{2}-\d{2}", cleaned) or "last_updated" in cleaned.lower():
        return ("", cleaned)

    # Regex for currency, percentages, multipliers, or speed metrics
    match = re.search(r"(\$\d+(?:\.\d+)?[kMB]?|\d+(?:\.\d+)?%|\b\d+(?:\.\d+)?x\b|<?\b\d+s\b|\$\d+k?)", cleaned, re.IGNORECASE)
    if match and match.group(0) and 2 <= len(match.group(0)) <= 10:
        num = match.group(0)
        # Reject standalone 4-digit years like 2026 or 2025
        if re.match(r"^\d{4}$", num):
            return ("", cleaned)

        label = cleaned.replace(num, "")
        label = re.sub(r"^[:\-\s\•\(\)]+", "", label)
        label = re.sub(r"[\(\)]+$", "", label).strip()
        if len(label) >= 3:
            return (num, label)

    return ("", cleaned)


def build_pitch_deck_pptx(slides_data: List[Dict[str, Any]], startup_name: str) -> bytes:
    """Builds a complete 13-slide enterprise 16:9 PPTX deck matching the Web UI design and returns raw bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors matching Dark Slate Card UI
    COLOR_BG = RGBColor(15, 22, 36)        # #0F1624 Deep Navy
    COLOR_CARD = RGBColor(14, 23, 42)      # #0E172A Dark Slate Card Surface
    COLOR_BORDER = RGBColor(30, 41, 59)    # #1E293B Slate Border
    COLOR_CYAN = RGBColor(56, 189, 248)    # #38BDF8 Electric Cyan
    COLOR_WHITE = RGBColor(255, 255, 255)  # Title White
    COLOR_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Muted Subtitle
    COLOR_BODY = RGBColor(226, 232, 240)   # #E2E8F0 Body Text
    COLOR_BADGE_BG = RGBColor(30, 41, 59) # #1E293B Badge Surface

    for slide_data in slides_data:
        s_num = slide_data.get("slide_number", 1)
        s_num_str = f"0{s_num}" if s_num < 10 else str(s_num)
        s_type = (slide_data.get("slide_type") or "SLIDE").upper()
        title = clean_bullet_text(slide_data.get("title"), max_len=80)
        subtitle = clean_bullet_text(slide_data.get("subtitle"), max_len=120)
        content = clean_bullet_text(slide_data.get("content"), max_len=200)
        raw_kps = slide_data.get("key_points", [])

        s = prs.slides.add_slide(blank_layout)

        # 1. Background Fill (16:9 Widescreen)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # 2A. Header Tag (Left)
        tx_box = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.35))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"SLIDE {s_num_str} / 13"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

        # 2B. Header Type Pill Badge (Right)
        badge_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.533), Inches(0.45), Inches(2.0), Inches(0.35))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = COLOR_BADGE_BG
        badge_box.line.color.rgb = COLOR_BORDER
        badge_box.line.width = Pt(1)
        tf_b = badge_box.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        p_b.text = s_type
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = RGBColor(203, 213, 225)

        # 3. Slide Main Title
        tx_box = s.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.7))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        # 4A. Slide Subtitle
        tx_box = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_MUTED

        # 4B. Slide Content Description
        if content and content != subtitle:
            tx_box_c = s.shapes.add_textbox(Inches(0.8), Inches(2.15), Inches(11.733), Inches(1.2))
            tf_c = tx_box_c.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = 0
            tf_c.margin_top = 0
            p_c = tf_c.paragraphs[0]
            p_c.text = content
            p_c.font.size = Pt(11.5)
            p_c.font.color.rgb = RGBColor(203, 213, 225)

        # Component 1: Cover Slide
        if s_num == 1 or slide_data.get("slide_type") == "cover":
            card_y = Inches(3.6)
            card_h = Inches(2.85)
            gap = Inches(0.25)
            card_w = Inches(2.745)  # (11.733 - 3*0.25) / 4

            for idx, kp in enumerate(raw_kps[:4]):
                col = idx % 4
                pos_x = Inches(0.8) + col * (card_w + gap)
                metric_num, metric_label = extract_metric(kp)

                card_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, card_y, card_w, card_h)
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = COLOR_CARD
                card_shape.line.color.rgb = COLOR_BORDER
                card_shape.line.width = Pt(1.5)

                tf_card = card_shape.text_frame
                tf_card.vertical_anchor = MSO_ANCHOR.TOP
                tf_card.word_wrap = True
                tf_card.margin_left = Inches(0.2)
                tf_card.margin_right = Inches(0.2)
                tf_card.margin_top = Inches(0.2)
                tf_card.margin_bottom = Inches(0.15)

                if metric_num:
                    p1 = tf_card.paragraphs[0]
                    p1.text = metric_num
                    p1.font.size = Pt(24)
                    p1.font.bold = True
                    p1.font.color.rgb = COLOR_CYAN

                    p2 = tf_card.add_paragraph()
                    p2.text = metric_label
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = COLOR_MUTED
                else:
                    p1 = tf_card.paragraphs[0]
                    clean_kp = clean_bullet_text(kp, max_len=130)
                    bullet_icon = "⚠️ " if "problem" in title.lower() or "risk" in title.lower() else "▸ "
                    if not clean_kp.startswith(("▸", "⚠️", "•", "✔", "➔")):
                        clean_kp = f"{bullet_icon}{clean_kp}"
                    p1.text = clean_kp
                    p1.font.size = Pt(11)
                    p1.font.color.rgb = COLOR_BODY

        # Component 2: Product Workflow Timeline Pipeline Slide (Slide 6)
        elif s_num == 6 or slide_data.get("slide_type") == "product_workflow" or slide_data.get("visual_type") == "three_step_flow":
            steps = [
                {"step": 1, "title": "Context Ingestion", "desc": "Securely ingests startup profile & validation metrics."},
                {"step": 2, "title": "AI Execution", "desc": "Runs multi-agent Zero-Trust compliance & risk scoring."},
                {"step": 3, "title": "Audit Output", "desc": "Generates investor-ready presentation decks in seconds."}
            ]

            if slide_data.get("visual_data") and isinstance(slide_data["visual_data"].get("steps"), list):
                extracted_steps = slide_data["visual_data"]["steps"]
                if len(extracted_steps) >= 3:
                    steps = [
                        {
                            "step": st.get("step") or st.get("step_number") or i+1,
                            "title": clean_bullet_text(st.get("title") or st.get("step_title") or f"Step {i+1}", max_len=40),
                            "desc": clean_bullet_text(st.get("desc") or st.get("description") or "", max_len=90)
                        }
                        for i, st in enumerate(extracted_steps[:3])
                    ]

            pipe_y = Inches(3.6)
            pipe_h = Inches(2.85)
            pipe_w = Inches(3.411)  # (11.733 - 2*0.75) / 3
            pipe_gap = Inches(0.75)

            for idx, st in enumerate(steps):
                pos_x = Inches(0.8) + idx * (pipe_w + pipe_gap)

                pipe_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pipe_y, pipe_w, pipe_h)
                pipe_card.fill.solid()
                pipe_card.fill.fore_color.rgb = COLOR_CARD
                pipe_card.line.color.rgb = COLOR_BORDER
                pipe_card.line.width = Pt(1.5)

                tf_p = pipe_card.text_frame
                tf_p.vertical_anchor = MSO_ANCHOR.TOP
                tf_p.word_wrap = True
                tf_p.margin_left = Inches(0.2)
                tf_p.margin_right = Inches(0.2)
                tf_p.margin_top = Inches(0.2)
                tf_p.margin_bottom = Inches(0.15)

                # Step Badge Number
                p_num = tf_p.paragraphs[0]
                p_num.text = f"STEP {st['step']}"
                p_num.font.size = Pt(11)
                p_num.font.bold = True
                p_num.font.color.rgb = COLOR_CYAN

                # Step Title
                p_title = tf_p.add_paragraph()
                p_title.text = st['title']
                p_title.font.size = Pt(14)
                p_title.font.bold = True
                p_title.font.color.rgb = COLOR_WHITE

                # Step Description
                p_desc = tf_p.add_paragraph()
                p_desc.text = st['desc']
                p_desc.font.size = Pt(11)
                p_desc.font.color.rgb = COLOR_MUTED

                # Connecting Chevron Arrow
                if idx < len(steps) - 1:
                    arr_box = s.shapes.add_textbox(pos_x + pipe_w, pipe_y + Inches(1.1), pipe_gap, Inches(0.6))
                    tf_arr = arr_box.text_frame
                    p_arr = tf_arr.paragraphs[0]
                    p_arr.alignment = PP_ALIGN.CENTER
                    p_arr.text = "➔"
                    p_arr.font.size = Pt(20)
                    p_arr.font.bold = True
                    p_arr.font.color.rgb = COLOR_CYAN

        # Component 3 & 4: Standard 4-Card & Metric Stat Component Slides
        else:
            card_y = Inches(3.6)
            card_h = Inches(2.85)
            gap = Inches(0.25)
            card_w = Inches(2.745)  # (11.733 - 3*0.25) / 4

            for idx, kp in enumerate(raw_kps[:4]):
                col = idx % 4
                pos_x = Inches(0.8) + col * (card_w + gap)
                metric_num, metric_label = extract_metric(kp)

                card_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, card_y, card_w, card_h)
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = COLOR_CARD
                card_shape.line.color.rgb = COLOR_BORDER
                card_shape.line.width = Pt(1.5)

                tf_card = card_shape.text_frame
                tf_card.vertical_anchor = MSO_ANCHOR.TOP
                tf_card.word_wrap = True
                tf_card.margin_left = Inches(0.2)
                tf_card.margin_right = Inches(0.2)
                tf_card.margin_top = Inches(0.2)
                tf_card.margin_bottom = Inches(0.15)

                if metric_num:
                    p1 = tf_card.paragraphs[0]
                    p1.text = metric_num
                    p1.font.size = Pt(24)
                    p1.font.bold = True
                    p1.font.color.rgb = COLOR_CYAN

                    p2 = tf_card.add_paragraph()
                    p2.text = metric_label
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = COLOR_MUTED
                else:
                    p1 = tf_card.paragraphs[0]
                    clean_kp = clean_bullet_text(kp, max_len=130)
                    bullet_icon = "⚠️ " if "problem" in title.lower() or "risk" in title.lower() else "▸ "
                    if not clean_kp.startswith(("▸", "⚠️", "•", "✔", "➔")):
                        clean_kp = f"{bullet_icon}{clean_kp}"
                    p1.text = clean_kp
                    p1.font.size = Pt(11)
                    p1.font.color.rgb = COLOR_BODY

        # 5. Footer Line & Text
        footer_box = s.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.4))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"VentureAI AI Pitch Deck Generator   |   {startup_name}"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = RGBColor(100, 116, 139)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
